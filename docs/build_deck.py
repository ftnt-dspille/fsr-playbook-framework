#!/usr/bin/env python3
"""FortiSOAR AI Investigation & Response connector — PM deep-dive deck (.pptx).

Run: .venv/bin/python docs/build_deck.py
Output: docs/FSR_Playbook_AI_Deepdive.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette --------------------------------------------------------------
BG      = RGBColor(0x0D, 0x11, 0x17)
PANEL   = RGBColor(0x16, 0x1C, 0x26)
PANEL2  = RGBColor(0x1E, 0x26, 0x33)
INK     = RGBColor(0xE6, 0xED, 0xF3)
MUTE    = RGBColor(0x9A, 0xA7, 0xB4)
ACCENT  = RGBColor(0x4C, 0x9A, 0xFF)
ACCENT2 = RGBColor(0x3F, 0xD9, 0x9B)
WARN    = RGBColor(0xFF, 0xB4, 0x54)
RED     = RGBColor(0xFF, 0x6B, 0x6B)
PURPLE  = RGBColor(0xB4, 0x8E, 0xFF)
CODEBG  = RGBColor(0x0A, 0x0E, 0x14)
LINE    = RGBColor(0x2A, 0x33, 0x40)

MONO = "Consolas"
SANS = "Calibri"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---- primitives -----------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def note(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip()


def _font(run, size, color, bold, font, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.name = font; run.font.italic = italic


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        p.space_before = Pt(0); p.line_spacing = line_spacing
        for seg in para:
            r = p.add_run(); r.text = seg[0]
            _font(r, seg[1], seg[2], seg[3], seg[4], seg[5] if len(seg) > 5 else False)
    return tb


def kicker(s, label):
    box(s, 0.6, 0.55, 0.12, 0.42, fill=ACCENT)
    text(s, 0.85, 0.5, 11, 0.5, [[(label.upper(), 13, ACCENT, True, SANS)]])


def title(s, t, sub=None):
    text(s, 0.6, 0.95, 12.1, 1.0, [[(t, 29, INK, True, SANS)]])
    if sub:
        text(s, 0.62, 1.6, 12.1, 0.5, [[(sub, 15, MUTE, False, SANS)]])


def footer(s, n):
    text(s, 0.6, 7.05, 9, 0.3,
         [[("FortiSOAR AI Investigation & Response  ·  the connector", 9, LINE, False, SANS)]])
    text(s, 12.2, 7.05, 0.9, 0.3, [[(str(n), 9, MUTE, False, SANS)]], align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, items, size=14, gap=7, color=INK, bullet_color=ACCENT):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            runs.append([("▸ ", size, bullet_color, True, SANS),
                         (head + "  ", size, color, True, SANS),
                         (body, size, MUTE, False, SANS)])
        else:
            runs.append([("▸ ", size, bullet_color, True, SANS),
                         (it, size, color, False, SANS)])
    text(s, x, y, w, 5, runs, space_after=gap, line_spacing=1.05)


def code(s, x, y, w, h, lines, size=12.5, title_txt=None):
    box(s, x, y, w, h, fill=CODEBG, line=LINE, line_w=1.0, radius=True)
    oy = y + 0.12
    if title_txt:
        text(s, x + 0.2, oy, w - 0.4, 0.3, [[(title_txt, 10.5, MUTE, True, MONO)]])
        oy += 0.34
    runs = [ln if isinstance(ln, list) else [(ln, size, INK, False, MONO)] for ln in lines]
    text(s, x + 0.22, oy, w - 0.44, h - (oy - y) - 0.1, runs, space_after=2, line_spacing=1.04)


def chip(s, x, y, label, color, w=None):
    w = w or (0.155 * len(label) + 0.3)
    box(s, x, y, w, 0.34, fill=PANEL2, line=color, line_w=1.0, radius=True)
    text(s, x, y + 0.02, w, 0.3, [[(label, 11, color, True, SANS)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


# ===========================================================================
# 1 — TITLE
# ===========================================================================
s = slide()
box(s, 0, 0, 13.333, 0.12, fill=ACCENT)
text(s, 0.9, 1.95, 11.7, 1.5,
     [[("The connector that turns FortiSOAR", 36, INK, True, SANS)],
      [("into an autonomous SOC analyst.", 36, ACCENT, True, SANS)]], line_spacing=1.05)
text(s, 0.92, 3.8, 11.2, 0.9,
     [[("AI investigates, hunts and enriches across your stack — then ", 17, INK, False, SANS),
       ("bottles that investigation into a repeatable playbook", 17, ACCENT2, True, SANS),
       (" any analyst can run, forever.", 17, INK, False, SANS)]],
     line_spacing=1.12)
cx = 0.92
for lbl, c in [("Investigate", ACCENT), ("Hunt", PURPLE), ("Enrich", ACCENT2),
               ("Contain (HITL)", WARN), ("→ Playbook", ACCENT)]:
    cx += chip(s, cx, 4.95, lbl, c) + 0.18
text(s, 0.92, 6.5, 11, 0.4,
     [[("In-platform connector (MCP agent)  ·  AngularJS widget  ·  live on FortiCloud SOAR 7.6.5", 13, MUTE, False, SANS)]])
note(s, """
Opening. This is one connector that ships inside FortiSOAR / FortiSOC. Frame for PMs: we are not adding a chat assistant — we are adding an analyst.
Three things to land in 30 seconds:
1) It actually does the investigation across the real stack, not a summary.
2) Any action that changes state is gated by a human.
3) The wow: the AI's investigation is automatically captured as a deterministic playbook — so we use AI to FIND the pattern once, then run it forever without AI. That's the line that should make a PM lean in.
Note 'live on FortiCloud SOAR 7.6.5' — this is real and running, not a concept.
""")


# ===========================================================================
# 2 — PROBLEM
# ===========================================================================
s = slide(); kicker(s, "The problem with 'AI for the SOC'")
title(s, "Most “AI in the SOC” is a summary box",
      "It reads the alert and writes a paragraph. It never touches your stack.")
cards = [
    ("Summarizes, doesn't investigate", RED,
     ["It restates the alert fields.", "It doesn't query FortiSIEM,",
      "doesn't pivot the indicators,", "doesn't pull the events that",
      "actually triggered the case."]),
    ("Hallucinated enrichment", WARN,
     ["“This IP is likely malicious”", "with no source behind it.",
      "No VirusTotal verdict, no", "Shodan exposure, no real",
      "reputation lookup you can cite."]),
    ("Can't be trusted to act", ACCENT,
     ["Containment means blocking,", "isolating, disabling — real",
      "state changes. A chatbot that", "might hallucinate must never",
      "fire those unsupervised."]),
]
x = 0.6
for head, c, body in cards:
    box(s, x, 2.2, 3.95, 3.7, fill=PANEL, line=LINE, radius=True)
    box(s, x, 2.2, 3.95, 0.1, fill=c)
    text(s, x + 0.25, 2.45, 3.5, 0.7, [[(head, 15.5, INK, True, SANS)]], line_spacing=1.0)
    text(s, x + 0.25, 3.35, 3.5, 2.4,
         [[(l, 12, MUTE, False, SANS)] for l in body], space_after=4, line_spacing=1.08)
    x += 4.18
text(s, 0.6, 6.25, 12.2, 0.6,
     [[("The fix isn't a smarter summary. ", 16, INK, True, SANS),
       ("It's wiring the model to the platform — and gating what it can do.", 16, ACCENT2, True, SANS)]])
note(s, """
Set up the contrast. Every vendor is shipping 'AI in the SOC' right now and 90% of it is a text box that summarizes the alert. PMs have seen ten of those demos this quarter.
Three failure modes — say them plainly:
- It summarizes but never investigates. It can't tell you anything the alert didn't already say.
- Its enrichment is made up. 'Probably malicious' with no source is worse than nothing — it can't be put in a ticket or trusted.
- It can't be trusted to take action, because if it can hallucinate a fact it can hallucinate a block.
Land the bottom line: the gap isn't model quality, it's plumbing and control. We solved both.
""")
footer(s, 2)


# ===========================================================================
# 3 — THESIS
# ===========================================================================
s = slide(); kicker(s, "What the connector is")
title(s, "It runs the investigation, not its mouth")
box(s, 0.6, 2.15, 12.1, 1.55, fill=PANEL, line=LINE, radius=True)
text(s, 0.95, 2.4, 11.5, 1.15,
     [[("A reasoning loop ", 18, INK, False, SANS),
       ("wired to your stack", 18, ACCENT, True, SANS),
       (": it pulls the triggering events from FortiSIEM / FortiAnalyzer, ", 18, INK, False, SANS),
       ("enriches every indicator", 18, ACCENT2, True, SANS),
       (" against real threat-intel sources, and ", 18, INK, False, SANS),
       ("stages containment for approval", 18, WARN, True, SANS),
       (" — every action tier-gated, logged, and reachable even on-prem.", 18, INK, False, SANS)]],
     line_spacing=1.15)
pairs = [
    ("Restates the alert", "Pulls the events that triggered it"),
    ("“Probably malicious”", "Cited verdicts from VT / Shodan / FortiGuard"),
    ("Suggests you go block it", "Stages the block; analyst approves; it runs"),
    ("Leaves you where you started", "Bottles the run into a reusable playbook"),
]
y = 4.05
for a, b in pairs:
    box(s, 0.6, y, 5.7, 0.62, fill=PANEL2, radius=True)
    text(s, 0.85, y + 0.05, 5.4, 0.5, [[("✗  " + a, 14, MUTE, False, SANS)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.45, y + 0.05, 0.5, 0.5, [[("→", 16, ACCENT, True, SANS)]], anchor=MSO_ANCHOR.MIDDLE)
    box(s, 7.0, y, 5.7, 0.62, fill=PANEL, line=ACCENT2, line_w=1.0, radius=True)
    text(s, 7.25, y + 0.05, 5.4, 0.5, [[("✓  " + b, 14, INK, True, SANS)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.7
note(s, """
This is the thesis slide — the one-liner is 'it runs the investigation, not its mouth.'
Walk the four rows left-to-right; each is a 'them vs us':
- A summary restates the alert; we pull the events that triggered it out of SIEM/FAZ.
- They say 'probably malicious'; we return cited verdicts you can paste into a ticket.
- They suggest you go block it; we stage the block, the analyst approves, and it actually executes — even against an isolated on-prem box.
- They leave you where you started; we hand you a reusable playbook.
The phrase 'reachable even on-prem' in the box is a hook for slide 11 — don't over-explain it here, just plant it.
""")
footer(s, 3)


# ===========================================================================
# 4 — WHERE IT LIVES (connector = engine)
# ===========================================================================
s = slide(); kicker(s, "The connector is the engine")
title(s, "The widget is a window. The connector does the work.",
      "All intelligence, tools, gating and trace-capture live in the in-platform connector")
box(s, 0.6, 2.25, 3.7, 1.5, fill=PANEL, line=ACCENT, line_w=1.0, radius=True)
text(s, 0.8, 2.4, 3.4, 0.4, [[("AngularJS widget", 15, INK, True, SANS)]])
text(s, 0.8, 2.82, 3.4, 0.9, [[("A window onto the agent. Renders cards, halts for approval, resumes.", 11.5, MUTE, False, SANS)]], line_spacing=1.05)
text(s, 4.35, 2.85, 0.5, 0.4, [[("→", 20, MUTE, True, SANS)]])
box(s, 4.85, 2.2, 3.7, 1.6, fill=PANEL, line=ACCENT2, line_w=1.75, radius=True)
text(s, 5.05, 2.33, 3.4, 0.4, [[("The connector (agent)", 15, ACCENT2, True, SANS)]])
text(s, 5.05, 2.74, 3.4, 1.0, [[("Tier-gated tool loop · HITL durability · remote execution · trace capture.", 11.5, MUTE, False, SANS)]], line_spacing=1.05)
text(s, 8.6, 2.85, 0.5, 0.4, [[("→", 20, MUTE, True, SANS)]])
box(s, 9.1, 2.25, 3.6, 1.5, fill=PANEL, line=WARN, line_w=1.0, radius=True)
text(s, 9.3, 2.4, 3.2, 0.4, [[("Your stack", 15, INK, True, SANS)]])
text(s, 9.3, 2.82, 3.2, 0.9, [[("FortiSIEM · FortiAnalyzer · EDR · firewall · 3rd-party SIEM · TI.", 11.5, MUTE, False, SANS)]], line_spacing=1.05)
box(s, 0.6, 4.15, 12.1, 2.45, fill=PANEL2, line=LINE, radius=True)
text(s, 0.85, 4.3, 11.5, 0.4, [[("What the connector can reach — by live discovery, never guessed", 14, INK, True, SANS)]])
grid = [("FortiSIEM", "siem_search_ip / host / user · events_for_incident", PURPLE),
        ("FortiAnalyzer", "faz_get_alerts · faz_search_ip · faz_raw_query", PURPLE),
        ("Threat intel", "VirusTotal · Shodan · AbuseIPDB · FortiGuard · URLScan", ACCENT2),
        ("Response", "firewall block · EDR isolate · identity disable", WARN),
        ("FortiSOAR data", "get_record (curated) · module search", ACCENT),
        ("714+ connectors", "any 3rd-party tool — find_connector / find_operation", MUTE)]
for i, (h, b, c) in enumerate(grid):
    gx = 0.85 + (i % 3) * 3.95
    gy = 4.8 + (i // 3) * 0.85
    box(s, gx, gy, 3.8, 0.72, fill=PANEL, line=LINE, radius=True)
    text(s, gx + 0.18, gy + 0.08, 3.5, 0.3, [[(h, 12.5, c, True, SANS)]])
    text(s, gx + 0.18, gy + 0.4, 3.5, 0.3, [[(b, 10, MUTE, False, MONO)]])
note(s, """
Important framing for this audience: the widget is NOT the product. The connector is. The widget is just a window — everything that matters (the reasoning loop, the tools, the safety gating, the trace capture) lives in the connector, server-side, inside FortiSOAR.
Why this matters to a PM: it means this rides the entire FortiSOAR distribution and licensing model. No new appliance, no new console, no agent to push to endpoints — it's a connector. Install it the way you install any connector.
The reach grid is the proof of 'wired to the stack': name a couple — FortiSIEM hunting tools, real TI sources, and the 714+ connector library that means 'any third-party tool' including non-Fortinet SIEMs. Tee up that last cell for the vendor-agnostic slide.
""")
footer(s, 4)


# ===========================================================================
# 5 — THE LOOP
# ===========================================================================
s = slide(); kicker(s, "The investigation loop")
title(s, "One alert → verdict → response → reusable playbook")
stepc = [("Ground", "normalize the alert; pull triggering events", ACCENT),
         ("Hunt", "pivot indicators across SIEM / FAZ", PURPLE),
         ("Enrich", "cited verdicts from TI sources", ACCENT2),
         ("Assess", "forced verdict + next action", INK),
         ("Contain", "stage action · analyst approves", WARN),
         ("Capture", "trace → re-runnable playbook", ACCENT)]
x = 0.6
for i, (h, b, c) in enumerate(stepc):
    box(s, x, 2.7, 1.85, 1.6, fill=PANEL, line=c, line_w=1.25, radius=True)
    box(s, x, 2.7, 1.85, 0.1, fill=c)
    text(s, x + 0.12, 2.92, 1.6, 0.4, [[(str(i+1) + ". " + h, 14, INK, True, SANS)]])
    text(s, x + 0.12, 3.32, 1.62, 0.95, [[(b, 11, MUTE, False, SANS)]], line_spacing=1.05)
    if i < 5:
        text(s, x + 1.82, 3.18, 0.32, 0.4, [[("→", 15, MUTE, True, SANS)]])
    x += 2.04
box(s, 0.6, 4.75, 12.1, 0.78, fill=PANEL2, line=LINE, radius=True)
text(s, 0.85, 4.86, 11.6, 0.6,
     [[("Every step is a tool call, recorded. ", 14.5, INK, True, SANS),
       ("Read-only steps run freely; any state-changing step suspends for approval before it fires.", 14.5, MUTE, False, SANS)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.6, 5.9, 12, 0.7,
     [[("The analyst never leaves the record, never writes a query — and ends with a verdict, an audit trail, and a playbook they can run next time.", 14, ACCENT2, False, SANS)]], line_spacing=1.1)
note(s, """
This is the map for the next several slides — point at it and say 'we'll drill into each box.'
The key idea is that the whole loop is made of recorded tool calls. That recording is what makes step 6 (Capture) possible — hold that thought, it pays off as the hero slide later.
Emphasize the bottom line for PMs: the analyst never leaves the FortiSOAR record and never writes a SIEM query by hand. That's the day-to-day time save. The verdict + audit trail + playbook is what they walk away with.
""")
footer(s, 5)


# ===========================================================================
# 6 — GROUNDING
# ===========================================================================
s = slide(); kicker(s, "Grounding")
title(s, "The agent starts from evidence, not the prompt",
      "The context is built before the model ever sees it")
layers = [
    ("L0 · Normalize", ACCENT,
     ["Collapse alert + case into one shape",
      "Extract IPs / hosts / users / hashes",
      "from first-class fields AND free text",
      "Surface MITRE tactics + severity",
      "Pull the FULL triggering events —",
      "the evidence a record view truncates"]),
    ("L1 · Classify", PURPLE,
     ["Detect the scenario:",
      "C2 / exfil · malware · mail-egress",
      "intrusion · defense-evasion",
      "",
      "Map → scenario-specific",
      "“opening moves” for the agent"]),
    ("L2 · Build prompt", ACCENT2,
     ["Inject a 'what we know' block",
      "Pre-fill entity context",
      "Point at first-class SIEM tools",
      "(not fiddly raw queries)",
      "Attach a verdict + containment",
      "checklist"]),
]
x = 0.6
for head, c, items in layers:
    box(s, x, 2.2, 3.92, 4.2, fill=PANEL, line=LINE, radius=True)
    box(s, x, 2.2, 3.92, 0.1, fill=c)
    text(s, x + 0.22, 2.42, 3.5, 0.4, [[(head, 16, INK, True, SANS)]])
    text(s, x + 0.22, 2.95, 3.55, 3.3,
         [[(it, 12.5, INK if it else MUTE, False, SANS)] for it in items],
         space_after=5, line_spacing=1.05)
    x += 4.18
text(s, 0.6, 6.6, 12.2, 0.5,
     [[("Result: ", 13.5, ACCENT2, True, SANS),
       ("the model reasons over real, structured evidence instead of mining a 100 KB truncated record live.", 13.5, INK, False, SANS)]])
note(s, """
This is the 'why it's accurate, not just confident' slide. Skeptics worry AI makes things up — the answer is we don't let it start from a blank prompt.
Three layers, plain language:
- Normalize: we extract the real indicators (IPs, hosts, users, hashes) and, critically, pull the FULL set of triggering events — the evidence the alert view usually truncates. The model sees the facts.
- Classify: we recognize the scenario (C2, malware, exfil…) and hand the agent proven opening moves, like an experienced analyst's playbook-in-the-head.
- Build prompt: we assemble a 'what we know' brief and point it at the right tools.
PM takeaway: accuracy is engineered upstream. Less hallucination, fewer wasted tokens, faster verdicts.
""")
footer(s, 6)


# ===========================================================================
# 7 — HUNTING
# ===========================================================================
s = slide(); kicker(s, "Hunting")
title(s, "First-class hunting across SIEM & FortiAnalyzer",
      "The agent pivots indicators in every source — with the right tool, not raw queries")
bullets(s, 0.6, 2.2, 6.0, [
    ("Source-matched", "SIEM alert → siem_* tools; FAZ alert → faz_* tools."),
    ("Indicator pivots", "siem_search_ip / host / user across a time window."),
    ("Incident events", "siem_events_for_incident pulls what triggered the case."),
    ("Escape hatches", "siem_raw_query / faz_raw_query for native queries."),
    ("Resilient transport", "modern JSON APIs, not the legacy query that 400s."),
    ("Goes wide", "independent lookups batched in one turn for speed."),
], size=13.5, gap=10)
code(s, 6.85, 2.2, 5.85, 4.4,
     [[("siem_search_ip(", 12.5, INK, False, MONO)],
      [('  "102.220.160.21",', 12.5, WARN, False, MONO)],
      [('  direction="dst", window="2h")', 12.5, ACCENT2, False, MONO)],
      [("", 5, INK, False, MONO)],
      [("→ events (digested):", 12, MUTE, True, MONO)],
      [("  ts  src→dst  bytes  svc  action", 11.5, INK, False, MONO)],
      [("", 5, INK, False, MONO)],
      [("# friendly fields → backend attrs:", 10.5, MUTE, False, MONO)],
      [("#   sourceip → srcIpAddr", 11, ACCENT, False, MONO)],
      [("#   dstipv4  → destIpAddr", 11, ACCENT, False, MONO)],
      [("", 5, INK, False, MONO)],
      [("# submit → poll → fetch, digested", 10.5, MUTE, False, MONO)],
      [("# to a compact result", 10.5, MUTE, False, MONO)]],
     title_txt="hunt FortiSIEM — no query language required")
note(s, """
Hunting is where analysts spend real time today, and it's specialist work — you need someone who knows the SIEM query language.
The point: the agent has purpose-built hunting tools. The analyst (or the AI) says 'show me everything talking to this IP in the last 2 hours' and we translate friendly fields to backend attributes, submit, poll, and digest the result.
Two credibility points for technical PMs:
- We matched the tool to the source (SIEM alert → SIEM tools, FAZ → FAZ tools) so it pivots in EVERY system, not just one.
- We hardened the transport — we use the modern JSON APIs because the legacy query path errors out on some versions. This is the kind of detail that says 'this was built by people who run it in production.'
""")
footer(s, 7)


# ===========================================================================
# 8 — ENRICHMENT
# ===========================================================================
s = slide(); kicker(s, "Enrichment")
title(s, "Cited verdicts, discovered live — never invented",
      "The connector surfaces real, configured, read-only intel ops")
bullets(s, 0.6, 2.2, 6.0, [
    ("Discovery, not guessing", "every op comes from the live connector store."),
    ("Ranked sources", "VirusTotal / FortiGuard first; Shodan / AbuseIPDB / URLScan next."),
    ("Read-only by tier", "only intel ops with no side effects can be picked."),
    ("Per-connector cap", "max 3 ops each, so one chatty connector can't crowd the slate."),
    ("Health-aware", "only configured, healthy connectors are offered."),
], size=13.5, gap=10)
code(s, 6.85, 2.2, 5.85, 2.5,
     [[("find_enrichment_actions(", 12.5, INK, False, MONO)],
      [('  target_type="ip")', 12.5, ACCENT2, False, MONO)],
      [("", 4, INK, False, MONO)],
      [("→ virustotal · ip_reputation", 11.5, ACCENT, False, MONO)],
      [("  shodan      · host_info", 11.5, ACCENT, False, MONO)],
      [("  abuseipdb   · check_ip", 11.5, ACCENT, False, MONO)]],
     title_txt="enrichment slate (real ops)")
box(s, 6.85, 4.9, 5.85, 1.7, fill=PANEL, line=LINE, radius=True)
text(s, 7.07, 5.05, 5.4, 0.4, [[("Source-aware record curation", 13.5, WARN, True, SANS)]])
text(s, 7.05, 5.45, 5.5, 1.1,
     [[("A hydrated FSR alert is ~100 KB. We drop audit noise, collapse picklists, and cap the payload", 12, MUTE, False, SANS),
       (" — keeping every indicator, losing the bloat.", 12, ACCENT2, False, SANS)]], line_spacing=1.08)
note(s, """
Enrichment is the antidote to the 'hallucinated verdict' fear from slide 2.
The connector doesn't have a hard-coded list of intel sources and it doesn't make up verdicts — it DISCOVERS which intel connectors you actually have configured and healthy, then uses those. Every verdict is attributable to a named op like virustotal/ip_reputation. You can put it in a ticket.
The ranking + per-connector cap detail shows this is curated, not a firehose — we lead with the best sources and don't let one chatty connector dominate.
The curation box is a nice efficiency proof: a hydrated alert is ~100KB of mostly audit noise; we shrink it to the indicators that matter. That's lower token cost and faster, sharper reasoning — a margin story a PM cares about.
""")
footer(s, 8)


# ===========================================================================
# 9 — CONTAINMENT / HITL
# ===========================================================================
s = slide(); kicker(s, "Containment")
title(s, "Real response actions — staged, never silent",
      "The connector discovers response ops; the analyst is always the trigger")
box(s, 0.6, 2.2, 5.95, 4.3, fill=PANEL, line=LINE, radius=True)
box(s, 0.6, 2.2, 5.95, 0.1, fill=WARN)
text(s, 0.85, 2.4, 5.4, 0.4, [[("Discovery → staging → approval", 16, INK, True, SANS)]])
bullets(s, 0.82, 2.95, 5.5, [
    ("Discover", "response verbs: block / quarantine / isolate / disable / revoke."),
    ("Target-matched", "by entity type — IP, host/endpoint, user, file."),
    ("Stage", "emit an action card with an editable arg preview."),
    ("Halt", "turn suspends; widget shows the card; nothing fires."),
    ("Approve", "analyst edits + confirms → resume carries the decision."),
    ("Execute", "the op runs with the approved args — and returns a result."),
], size=12.5, gap=9, bullet_color=WARN)
code(s, 6.78, 2.2, 5.95, 2.65,
     [[("emit_action_card(", 12, INK, False, MONO)],
      [('  connector="fortios",', 12, ACCENT, False, MONO)],
      [('  operation="block_ip",', 12, WARN, True, MONO)],
      [('  args={"ip": "102.220.160.21"},', 12, INK, False, MONO)],
      [('  editable_fields=["ip"])', 12, ACCENT2, False, MONO)],
      [("", 4, INK, False, MONO)],
      [("→ awaiting_action_card (halts)", 11, MUTE, False, MONO)]],
     title_txt="stage — does not execute")
box(s, 6.78, 5.05, 5.95, 1.45, fill=PANEL, line=ACCENT2, line_w=1.0, radius=True)
text(s, 7.0, 5.2, 5.5, 0.4, [[("Never a dead end", 13.5, ACCENT2, True, SANS)]])
text(s, 6.98, 5.6, 5.55, 0.85,
     [[("No response connector configured? ", 12, INK, False, SANS),
       ("a capability-gap card names the missing op, the connector to enable, and a resume button.", 12, MUTE, False, SANS)]], line_spacing=1.05)
note(s, """
This is the trust slide. The #1 objection to 'AI that acts' is 'what if it does something destructive.' Our answer is structural, not a promise: the AI can never silently fire a state-changing action.
Walk the flow: it DISCOVERS the response options, STAGES one as an editable card, and then HALTS. Nothing executes until the analyst edits the args and approves. The approval is durable — it survives across sessions, so an action staged now can be approved later.
'Never a dead end' is a deliberate product principle: if you don't have a block connector, it doesn't shrug — it tells you exactly which connector to enable and gives you a button to resume. The AI always hands the analyst a path forward.
Tee up slide 11: '…and approve doesn't just mean the API is reachable — watch what happens when the target is on-prem.'
""")
footer(s, 9)


# ===========================================================================
# 10 — TIER MODEL
# ===========================================================================
s = slide(); kicker(s, "The safety model")
title(s, "Every tool carries a tier — the gate is structural",
      "Approval isn't a prompt instruction. It's dispatch logic in the connector.")
tiers = [
    ("0", "Local only", "discovery, record lookups, card emission", ACCENT2, "auto"),
    ("1", "Read-only FSR / SIEM", "get_record, siem_*, faz_* hunting", ACCENT2, "auto"),
    ("2", "Read-only external", "VirusTotal / Shodan / AbuseIPDB enrichment", ACCENT, "auto"),
    ("3", "FSR data mutation", "update record, manage assets", WARN, "approval"),
    ("4", "Third-party state change", "block IP, isolate host, disable user", RED, "approval"),
]
y = 2.35
for t, name, ex, c, mode in tiers:
    box(s, 0.6, y, 12.1, 0.78, fill=PANEL, line=LINE, radius=True)
    box(s, 0.6, y, 0.7, 0.78, fill=PANEL2)
    text(s, 0.6, y + 0.18, 0.7, 0.5, [[(t, 22, c, True, SANS)]], align=PP_ALIGN.CENTER)
    text(s, 1.5, y + 0.13, 4.0, 0.5, [[(name, 15, INK, True, SANS)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.6, y + 0.13, 5.2, 0.5, [[(ex, 12.5, MUTE, False, SANS)]], anchor=MSO_ANCHOR.MIDDLE)
    mc = ACCENT2 if mode == "auto" else WARN
    box(s, 11.05, y + 0.19, 1.45, 0.4, fill=None, line=mc, line_w=1.25, radius=True)
    text(s, 11.05, y + 0.21, 1.45, 0.36,
         [[(("▶ runs" if mode == "auto" else "⏸ approval"), 11.5, mc, True, SANS)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.86
text(s, 0.6, 6.75, 12.2, 0.4,
     [[("A tier-≥3 call doesn't ask nicely — it ", 13, INK, True, SANS),
       ("suspends and cannot dispatch until a human decision is carried back in.", 13, MUTE, False, SANS)]])
note(s, """
Back up the trust claim with the mechanism. Every single tool the AI can call is tagged with a tier 0–4. Read-only things (discovery, hunting, enrichment) run automatically. Anything that changes state — tier 3 and 4 — physically cannot dispatch without a human decision.
The line that matters: this is dispatch logic, not a prompt instruction. We're not asking the model nicely to behave. The code refuses to execute the call. Even if the model were jailbroken or wrong, the gate holds.
For a PM thinking about enterprise sales and risk/compliance: this is the slide your security-conscious customers will want. It's auditable, it's deterministic, and it maps cleanly to 'analyst-in-the-loop' controls they already require.
""")
footer(s, 10)


# ===========================================================================
# 11 — REMOTE ACTIONS (HERO)
# ===========================================================================
s = slide(); kicker(s, "Wow factor · reach")
title(s, "Automate what the SaaS can't even reach",
      "FortiSOC is cloud. Your highest-value response targets are on-prem and isolated.")
box(s, 0.6, 2.2, 5.9, 2.0, fill=PANEL, line=LINE, radius=True)
text(s, 0.85, 2.35, 5.4, 0.4, [[("The reachability gap", 15, RED, True, SANS)]])
text(s, 0.85, 2.78, 5.45, 1.35,
     [[("The actions that matter most live where the cloud can't go: an ", 12.5, INK, False, SANS),
       ("on-prem AD server", 12.5, WARN, True, SANS),
       (", an internal firewall, an air-gapped segment. A SaaS analyst that can't reach them can investigate — but can't ", 12.5, INK, False, SANS),
       ("respond", 12.5, RED, True, SANS),
       (".", 12.5, INK, False, SANS)]], line_spacing=1.1)
box(s, 0.6, 4.4, 5.9, 2.2, fill=PANEL, line=ACCENT2, line_w=1.5, radius=True)
text(s, 0.85, 4.55, 5.4, 0.4, [[("The connector closes it", 15, ACCENT2, True, SANS)]])
bullets(s, 0.82, 5.0, 5.5, [
    "Drives actions through the on-prem FSR agent",
    "Executes locally — disable the AD account, block on the internal firewall",
    "And gets the real result back to the cloud",
    "Live-verified on FortiCloud SOAR 7.6.5",
], size=12, gap=8, bullet_color=ACCENT2)
code(s, 6.75, 2.2, 5.95, 4.4,
     [[("# target is behind an on-prem agent", 10.5, MUTE, False, MONO)],
      [("# (fire-and-forget by design)", 10.5, MUTE, False, MONO)],
      [("if config in agent_bound_configs:", 11.5, ACCENT, False, MONO)],
      [("", 4, INK, False, MONO)],
      [("  # run it as a 1-shot workflow", 10.5, MUTE, False, MONO)],
      [("  [ disable_ad_user ]", 11.5, WARN, True, MONO)],
      [("  [ force-fail to persist result ]", 11.5, INK, False, MONO)],
      [("", 4, INK, False, MONO)],
      [("  push · priority:High · trigger", 11.5, INK, False, MONO)],
      [("  poll → read the step result", 11.5, ACCENT, False, MONO)],
      [("  return {ok, data} to the cloud", 11.5, ACCENT2, False, MONO)],
      [("", 4, INK, False, MONO)],
      [("# the op runs exactly once", 10.5, ACCENT2, False, MONO)]],
     title_txt="remote execution — the hard part, solved")
note(s, """
This is a hero slide — slow down.
The strategic point: FortiSOC is a SaaS. The single biggest limitation of cloud security tooling is that it can't reach into the customer's network. And the most valuable response actions are exactly there — disabling a compromised account on an on-prem Active Directory server, blocking on an internal firewall, isolating a host on an air-gapped segment. A cloud analyst that can investigate but can't respond on-prem is only half a product.
The connector closes that gap. It drives the action through the on-prem FortiSOAR agent, executes it locally, and — this is the technically hard part — gets the real result back to the cloud. Agent execution is fire-and-forget by design (you fire the action and the platform never hands you the result), so we run it as a one-shot workflow engineered to capture and return the outcome. The op runs exactly once. No orphaned actions, no double-blocks.
Why PMs should care: a huge number of customers use FortiSOAR as their automation backbone precisely because it reaches everywhere. This makes the AI layer reach everywhere too. 'Automate anything, anywhere — including the stuff your SaaS competitors physically can't touch.'
Live-verified — say it.
""")
footer(s, 11)


# ===========================================================================
# 12 — GUARDRAILS
# ===========================================================================
s = slide(); kicker(s, "Guardrails")
title(s, "Determinism comes from code, not vibes",
      "Each guardrail is a mechanism in the loop with a unit test behind it")
guards = [
    ("Forced assessment", "Tool-only turn with no closing text → one forced summary round. Every investigation ends with a verdict + next action."),
    ("Repeated-error guard", "An identical call that already failed won't re-run — no loops, no budget burn."),
    ("Error classification", "A rejected request ('fix your params') is told apart from a real outage ('check connectivity')."),
    ("Low-signal gate", "'hi' / 'test' won't trigger an autonomous investigation. Intent is classified first."),
    ("Conditional-required", "A branch choice that activates nested required fields is validated — no silently-missing inputs."),
    ("Healthcheck cache", "Connector health is cached (hours when healthy), so the agent isn't re-probing the stack."),
]
for i, (head, desc) in enumerate(guards):
    bx = 0.6 + (i % 2) * 6.15
    by = 2.25 + (i // 2) * 1.42
    box(s, bx, by, 5.95, 1.28, fill=PANEL, line=LINE, radius=True)
    text(s, bx + 0.22, by + 0.14, 5.5, 0.35, [[(head, 14.5, INK, True, SANS)]])
    text(s, bx + 0.22, by + 0.5, 5.55, 0.7, [[(desc, 11.5, MUTE, False, SANS)]], line_spacing=1.04)
text(s, 0.6, 6.65, 12.2, 0.4,
     [[("Six of these ship today, ", 13, ACCENT2, True, SANS),
       ("each with its own unit test — the loop behaves the same way on run #1 and run #1000.", 13, MUTE, False, SANS)]])
note(s, """
Quick slide — don't read all six. The message is meta: 'this is engineered, not vibes.' Pick two examples:
- Forced assessment: the AI can't trail off after running tools — we force it to land a verdict and a next action every time. No 'here's some data, good luck.'
- Repeated-error guard: it won't get stuck retrying the same failing call. No runaway loops, no surprise bills.
Then land the closer: every one of these has a unit test, so the loop behaves identically on run #1 and run #1000. That consistency is exactly what you need before you trust automation — and it's what makes the playbook-capture story credible.
""")
footer(s, 12)


# ===========================================================================
# 13 — THE WOW: BOTTLE THE GENIUS
# ===========================================================================
s = slide(); kicker(s, "The wow factor")
title(s, "Bottle the moment of genius",
      "Use AI to find the pattern once. Run it as a playbook forever — no AI in the loop.")
# left: the insight
box(s, 0.6, 2.2, 5.95, 4.35, fill=PANEL, line=ACCENT2, line_w=1.5, radius=True)
text(s, 0.85, 2.38, 5.5, 0.45, [[("Why this changes the game", 16, ACCENT2, True, SANS)]])
bullets(s, 0.82, 2.95, 5.55, [
    ("AI explores faster than people", "tries query patterns, pivots, experiments with API calls in seconds."),
    ("But you shouldn't run it live forever", "cost, latency, and non-determinism don't belong in steady-state ops."),
    ("So we capture the winning run", "every recorded tool call becomes a playbook step — automatically."),
    ("The exploration happens once", "the playbook runs consistently, free, every time, by any analyst."),
], size=12.5, gap=11, bullet_color=ACCENT2)
# right: before/after
box(s, 6.78, 2.2, 5.95, 2.05, fill=PANEL, line=LINE, radius=True)
text(s, 7.0, 2.35, 5.5, 0.4, [[("Discover the pattern (AI, once)", 13, ACCENT, True, SANS)]])
text(s, 7.0, 2.75, 5.5, 1.4,
     [[("The agent investigates a novel threat — hunts in SIEM, enriches the indicators, confirms the verdict, stages the block. A real ", 11.5, MUTE, False, SANS),
       ("act of investigative reasoning.", 11.5, INK, True, SANS)]], line_spacing=1.08)
text(s, 9.4, 4.35, 0.6, 0.4, [[("↓", 22, ACCENT2, True, SANS)]])
box(s, 6.78, 4.85, 5.95, 1.7, fill=PANEL, line=ACCENT2, line_w=1.25, radius=True)
text(s, 7.0, 5.0, 5.5, 0.4, [[("Run the pattern (anyone, forever)", 13, ACCENT2, True, SANS)]])
text(s, 7.0, 5.4, 5.5, 1.1,
     [[("That exact investigation is compiled into a deterministic playbook. Wiring derived from the ", 11.5, MUTE, False, SANS),
       ("real values it observed", 11.5, INK, True, SANS),
       (" — not guessed. Verified, then offered to save.", 11.5, MUTE, False, SANS)]], line_spacing=1.08)
note(s, """
THE slide. If they remember one thing, make it this. Slow down and tell it as a story.
The insight: AI's real superpower in the SOC isn't answering questions — it's that it can EXPLORE faster than any human. It'll try ten query patterns, pivot across five indicators, experiment with API calls, in seconds. That's genuine investigative reasoning, compressed.
But here's the trap everyone else falls into: they then rely on the AI to RUN that investigation every single time. That's expensive, it's slow, and it's non-deterministic — you can't build a SOC's steady-state operations on 'the model will probably do the same thing again.'
Our move: we capture the winning run. Because every step of the investigation was a recorded tool call, we automatically compile it into a deterministic playbook — and the wiring between steps is derived from the REAL values the agent observed, not guessed. So the moment of genius happens once, with AI. Then it's bottled. Any L1 analyst can run that playbook on the next similar alert — instantly, consistently, at zero model cost, with no AI in the loop.
The soundbite: 'AI to discover the pattern. A playbook to run it. We bottle the moment of genius so it's not a one-time event — it's institutional knowledge any analyst can pour out on demand.'
This is also a moat story: every investigation your analysts run makes your playbook library richer. The product gets more valuable the more it's used.
""")
footer(s, 13)


# ===========================================================================
# 14 — WALKTHROUGH
# ===========================================================================
s = slide(); kicker(s, "End to end · C2 alert")
title(s, "wendy.smith / smithDesktop → 102.220.160.21 (NG)",
      "7ogger.exe · Exfiltration over C2 — what the analyst actually sees")
steps = [
    ("Open the alert", "the chat is already grounded in the normalized indicators — no prompt to write"),
    ("Pull the triggering events", "the agent hunts the C2 address in SIEM and confirms the beaconing"),
    ("Enrich the indicators", "VirusTotal on the malware hash; reputation on the destination IP — cited"),
    ("Verdict, forced", "“confirmed C2 / exfil, high severity — recommend blocking the destination”"),
    ("Stage containment", "an action card to block the IP — analyst reviews args and approves"),
    ("Bottle it", "the whole run is offered back as a reusable response playbook"),
]
y = 2.25
for i, (h, b) in enumerate(steps):
    box(s, 0.6, y, 0.55, 0.72, fill=PANEL2, line=ACCENT, radius=True)
    text(s, 0.6, y + 0.14, 0.55, 0.4, [[(chr(65+i), 16, ACCENT, True, SANS)]], align=PP_ALIGN.CENTER)
    box(s, 1.35, y, 11.35, 0.72, fill=PANEL, line=LINE, radius=True)
    text(s, 1.6, y + 0.1, 11, 0.35, [[(h, 14.5, INK, True, SANS)]])
    text(s, 1.6, y + 0.42, 11, 0.3, [[(b, 12, MUTE, False, SANS)]])
    y += 0.8
note(s, """
Make it concrete with a real scenario. This is a command-and-control / data-exfiltration case: a user's desktop is beaconing out to a Nigerian IP, with a suspicious binary (7ogger.exe) involved.
Walk the six steps as the analyst's lived experience — emphasize what they DIDN'T have to do:
- They open the alert. They don't write a prompt; it's already grounded.
- They don't write a SIEM query; the agent hunts and confirms the beaconing.
- They don't tab out to VirusTotal; enrichment comes back cited.
- They get a forced, decisive verdict — not a wishy-washy summary.
- They approve one block. That's the only click that matters, and it's theirs.
- And they walk away with a playbook for the next time this pattern shows up.
This is the 'day in the life' that makes the abstract slides real. If you have a recorded demo, this is where you'd play it.
""")
footer(s, 14)


# ===========================================================================
# 15 — VENDOR AGNOSTIC
# ===========================================================================
s = slide(); kicker(s, "Reach · any vendor")
title(s, "Same superpowers on anyone's stack",
      "FortiSOAR has integrated 3rd-party SIEMs and tools natively since day one")
box(s, 0.6, 2.2, 5.95, 2.0, fill=PANEL, line=LINE, radius=True)
text(s, 0.85, 2.35, 5.5, 0.4, [[("Not locked to Fortinet telemetry", 15, ACCENT, True, SANS)]])
text(s, 0.85, 2.78, 5.5, 1.35,
     [[("The connector enriches, hunts and contains through ", 12.5, INK, False, SANS),
       ("any of 714+ connectors", 12.5, ACCENT2, True, SANS),
       (" — third-party SIEMs, EDRs, firewalls, identity and TI. Same investigation quality whatever the customer runs.", 12.5, INK, False, SANS)]], line_spacing=1.1)
box(s, 0.6, 4.35, 5.95, 2.2, fill=PANEL, line=ACCENT2, line_w=1.25, radius=True)
text(s, 0.85, 4.5, 5.5, 0.4, [[("What FortiSOC actually adds", 15, ACCENT2, True, SANS)]])
text(s, 0.85, 4.95, 5.5, 1.5,
     [[("The integration was always there. FortiSOC ", 12.5, INK, False, SANS),
       ("unifies Fortinet logging, incident management and response under one product", 12.5, INK, True, SANS),
       (" — and this AI layer works across all of it, and everything it connects to.", 12.5, MUTE, False, SANS)]], line_spacing=1.12)
# right column: source-agnostic graphic
labels = [("3rd-party SIEM", PURPLE), ("Any EDR", WARN), ("Any firewall", WARN),
          ("Identity / IAM", ACCENT), ("Threat intel", ACCENT2), ("Cloud / SaaS", ACCENT)]
text(s, 6.85, 2.2, 5.85, 0.4, [[("One AI layer, every source", 14, INK, True, SANS)]])
for i, (lbl, c) in enumerate(labels):
    gx = 6.85 + (i % 2) * 3.0
    gy = 2.75 + (i // 2) * 1.25
    box(s, gx, gy, 2.85, 1.05, fill=PANEL, line=c, line_w=1.0, radius=True)
    text(s, gx, gy + 0.32, 2.85, 0.4, [[(lbl, 13, INK, True, SANS)]], align=PP_ALIGN.CENTER)
note(s, """
This neutralizes the 'we're an all-Fortinet shop / we're NOT a Fortinet shop' objection in both directions.
Key fact PMs should internalize and repeat: FortiSOAR has integrated with third-party SIEMs and security tools natively since it was created. That 714+ connector library isn't new — it's a decade-plus of integration work. This AI layer rides on ALL of it. So we can deliver the same enrichment, hunting and containment quality whether the customer runs Splunk, CrowdStrike, Palo Alto, whatever.
Then the honest, sharp positioning of FortiSOC: the integration breadth was always there. What FortiSOC changes is the GO-TO-MARKET — it unifies Fortinet logging, incident management, and response into one product we can sell as a unit. The AI layer makes that unified product feel like a single intelligent system rather than a bundle.
For a PM building the pitch: this is both a competitive moat (breadth) and a packaging story (unification). Use whichever the room needs.
""")
footer(s, 15)


# ===========================================================================
# 16 — ROADMAP
# ===========================================================================
s = slide(); kicker(s, "Where this goes next")
title(s, "From assistant to autonomous tier-1 SOC")
road = [
    ("AI investigation queue", PURPLE,
     ["A live watchtower of every AI",
      "investigation in flight. One",
      "central place to watch, approve,",
      "and guardrail — approvals surface",
      "to the analyst as they're needed."]),
    ("Auto-triage on case creation", ACCENT,
     ["Fires the moment a case is created.",
      "Runs full level-1 triage unattended,",
      "then escalates to a human with the",
      "findings already attached — the",
      "analyst starts at the verdict."]),
    ("False-positive auto-close + tuning", ACCENT2,
     ["For noisy escalation rules: confirm",
      "the FP, close the alert and case —",
      "and pinpoint exactly where to tune",
      "the FAZ / SIEM rule that fired,",
      "via an MCP link to the source."]),
]
x = 0.6
for head, c, items in road:
    box(s, x, 2.25, 3.92, 4.2, fill=PANEL, line=LINE, radius=True)
    box(s, x, 2.25, 3.92, 0.1, fill=c)
    text(s, x + 0.22, 2.48, 3.5, 0.75, [[(head, 15.5, INK, True, SANS)]], line_spacing=1.0)
    text(s, x + 0.22, 3.45, 3.55, 3.0,
         [[(it, 12.5, MUTE, False, SANS)] for it in items], space_after=6, line_spacing=1.08)
    x += 4.18
note(s, """
Close the vision. Today it's an assistant the analyst drives. The roadmap is autonomy with oversight.
1) AI investigation queue: as you scale to many concurrent AI investigations, you need a watchtower — one place to see everything in flight and a single, central spot where approvals and guardrails surface. This is the operational control plane for AI in the SOC, and it's a natural next product surface.
2) Auto-triage on case creation: flip it from analyst-initiated to event-driven. The moment a case is created, the connector runs full level-1 triage on its own, then escalates to a human WITH the findings attached. Your analysts stop starting from zero — they start at the verdict. This is the headline efficiency play: collapse the L1 tier.
3) False-positive auto-close + rule tuning: for noisy rules, it confirms the false positive, closes the alert and case, AND tells you precisely where to tune the FAZ/SIEM rule that fired — closing the loop back to the detection source via MCP. This attacks alert fatigue at the root, not just the symptom.
Tie it together: each of these is the same connector, same safety model, pointed at a bigger slice of the workflow. The destination is an autonomous tier-1 SOC that humans supervise instead of staff.
""")
footer(s, 16)


# ===========================================================================
# 17 — PROOF POINTS
# ===========================================================================
s = slide(); kicker(s, "Why this isn't fluff")
title(s, "Proof points, not promises")
proofs = [
    ("It touches the real stack", "FortiSIEM + FortiAnalyzer hunting, cited TI enrichment, firewall/EDR/identity response — all by live discovery.", ACCENT),
    ("It reaches on-prem & isolated", "Drives and returns results from agent-bound targets the cloud can't address directly. Live-proofed.", PURPLE),
    ("Action is gated by code", "State-changing ops structurally suspend for approval. The model cannot fire containment alone.", WARN),
    ("It compounds", "Every investigation can become a deterministic playbook — the library grows as analysts work.", ACCENT2),
]
for i, (h, b, c) in enumerate(proofs):
    bx = 0.6 + (i % 2) * 6.15
    by = 2.3 + (i // 2) * 2.05
    box(s, bx, by, 5.95, 1.85, fill=PANEL, line=LINE, radius=True)
    box(s, bx, by, 0.1, 1.85, fill=c)
    text(s, bx + 0.3, by + 0.22, 5.5, 0.5, [[(h, 17, INK, True, SANS)]])
    text(s, bx + 0.3, by + 0.78, 5.45, 1.0, [[(b, 13, MUTE, False, SANS)]], line_spacing=1.1)
note(s, """
Recap the defensible claims — four, one per quadrant:
- It touches the real stack (not a summary).
- It reaches on-prem and isolated targets (the SaaS-can't story) — and it's live-proofed, not theoretical.
- Action is gated by code (the trust story) — structural, not a prompt.
- It compounds — this is the business kicker. Unlike a feature that delivers fixed value, this gets MORE valuable the more it's used, because every investigation can be bottled into the playbook library.
End on 'compounds' — it's the line that turns a feature into a platform investment.
""")
footer(s, 17)


# ===========================================================================
# 18 — HOW IT SHIPS
# ===========================================================================
s = slide(); kicker(s, "How it ships")
title(s, "It's a connector — so it ships like one",
      "No new appliance, no endpoint agent to deploy, no separate console")
comps = [
    ("In-platform connector", "the engine", "the agent loop, tools, gating, remote execution, trace capture", ACCENT2),
    ("AngularJS widget", "the surface", "chat drawer on the record · renders cards · halt-and-resume", ACCENT),
    ("fsr_playbooks", "the shared core", "triage pipeline · tier model · compiler · reference store", WARN),
]
x = 0.6
for h, tag, b, c in comps:
    box(s, x, 2.3, 3.9, 2.3, fill=PANEL, line=c, line_w=1.25, radius=True)
    text(s, x + 0.22, 2.5, 3.5, 0.4, [[(h, 16, INK, True, SANS)]])
    text(s, x + 0.22, 2.92, 3.5, 0.3, [[(tag, 11.5, c, True, SANS)]])
    text(s, x + 0.22, 3.35, 3.5, 1.1, [[(b, 12.5, MUTE, False, SANS)]], line_spacing=1.12)
    x += 4.18
box(s, 0.6, 4.85, 12.1, 1.75, fill=PANEL2, line=LINE, radius=True)
text(s, 0.85, 5.0, 11.5, 0.4, [[("Already real today", 14, INK, True, SANS)]])
facts = [("Live", "FortiCloud SOAR 7.6.5"),
         ("Verified", "guardrails ship with unit tests"),
         ("Remote", "agent-bound execution proven"),
         ("Vendor-agnostic", "714+ connectors")]
cx = 0.85
for k, d in facts:
    box(s, cx, 5.45, 2.88, 0.95, fill=PANEL, line=LINE, radius=True)
    text(s, cx + 0.15, 5.55, 2.6, 0.3, [[(k, 12.5, ACCENT2, True, SANS)]])
    text(s, cx + 0.15, 5.9, 2.6, 0.5, [[(d, 10.5, MUTE, False, SANS)]], line_spacing=1.0)
    cx += 3.0
note(s, """
De-risk the 'how hard is this to adopt' question. It's a connector. It installs through the same mechanism as every other FortiSOAR connector — no new appliance, no endpoint agent rollout, no separate console to stand up or train on. The widget just surfaces it on the record the analyst already uses.
Reinforce the components: the connector is the engine, the widget is the surface, and fsr_playbooks is the shared brain that both the connector and the widget run.
The 'already real today' strip is your credibility anchor: it's live on a real build, the guardrails have tests, remote agent execution is proven, and it works across 700+ connectors. This is not a prototype slideware pitch — it's running.
""")
footer(s, 18)


# ===========================================================================
# 19 — CLOSE
# ===========================================================================
s = slide()
box(s, 0, 0, 13.333, 0.12, fill=ACCENT)
text(s, 0.9, 1.8, 11.6, 1.3,
     [[("AI finds the pattern.", 34, INK, True, SANS)],
      [("A playbook runs it forever.", 34, ACCENT2, True, SANS)]], line_spacing=1.1)
takeaways = [
    "It investigates across your real stack — Fortinet and any third-party tool or SIEM.",
    "It reaches on-prem and isolated targets the cloud alone can't — and returns the result.",
    "State-changing actions are tier-gated and analyst-approved — never silent.",
    "Every investigation can be bottled into a deterministic playbook — value compounds.",
]
y = 3.5
for t in takeaways:
    text(s, 0.95, y, 11.6, 0.5, [[("→  ", 16, ACCENT, True, SANS), (t, 16, INK, False, SANS)]])
    y += 0.62
text(s, 0.95, 6.45, 11, 0.4,
     [[("In-platform connector · AngularJS widget · live on FortiCloud SOAR 7.6.5", 12, MUTE, False, SANS)]])
note(s, """
Closing — land the slogan: 'AI finds the pattern; a playbook runs it forever.' That single sentence is the whole pitch.
Recap the four pillars (stack reach, on-prem reach, gated action, compounding value), then stop talking. Ask for what you came for — pilot customers, eng investment, a design-partner, whatever the ask is.
If you get one question, expect 'how do you stop it doing something dumb' — point back to the tier model (slide 10). The second most likely is 'does it work with non-Fortinet tools' — slide 15.
""")


prs.save("docs/FSR_Playbook_AI_Deepdive.pptx")
print("saved docs/FSR_Playbook_AI_Deepdive.pptx ·", len(prs.slides._sldIdLst), "slides")
