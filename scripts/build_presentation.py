"""Generate FSRPlaybookYaml presentation deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0F, 0x1E, 0x3D)
ACCENT = RGBColor(0xE8, 0x3E, 0x3E)
GREY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT = RGBColor(0xF4, 0xF4, 0xF6)

SLIDES = [
    {
        "title": "Agent-Authored Playbooks for FortiSOAR",
        "subtitle": "A YAML IR + a deterministic compiler + an MCP toolbelt — any LLM can drive it",
        "kind": "title",
    },
    {
        "title": "The pain we keep hitting",
        "bullets": [
            "Playbook Designer is click-heavy, opaque, and not diff-able",
            "Iteration loop is slow: edit → save → run → squint at run-log → repeat",
            "Field-shape bugs (.records vs .data, Jinja typos) only surface at runtime",
            "No version control, no PR review, no portability between environments",
            "Talented automation engineers spend their day clicking, not designing",
        ],
    },
    {
        "title": "The product principle",
        "bullets": [
            "\"Producing valid YAML is the start of the job, not the end.\"",
            "A playbook is useless if it doesn't run on a real FSR and produce the asked-for outcome",
            "So the system has to prove playbooks work, not just compile",
            "Everything that follows serves that one principle",
        ],
    },
    {
        "title": "Why now: agents are good enough",
        "bullets": [
            "Frontier LLMs handle structured authoring well — given the right tools",
            "Agents need: a queryable spec, a deterministic compiler, fast feedback",
            "No fine-tuning required — tool use plus good cheatsheets is sufficient",
            "Bottleneck moved from \"can the model write YAML\" to \"can it verify what it wrote\"",
        ],
    },
    {
        "title": "The core idea",
        "bullets": [
            "One simplified YAML IR. Three authoring surfaces (CLI, agent-via-MCP, future visual editor)",
            "All three compile through the same pipeline to native FSR WorkflowCollection JSON",
            "FortiSOAR stays the execution engine — we never replace it, we author into it",
            "Same IR across surfaces → output is interoperable",
        ],
        "diagram": "core",
    },
    {
        "title": "Why YAML is the right IR",
        "bullets": [
            "Human-readable, diff-able, review-able in a PR",
            "Trivially round-trippable: pull from FSR → YAML → edit → push",
            "Easy for an LLM to emit and reason about",
            "No vendor JSON ceremony (UUIDs, port coords, layout metadata) leaks into authoring",
        ],
    },
    {
        "title": "Design 1 — SQLite-first reference store",
        "bullets": [
            "714 connectors, 6,773 ops, 26K params, 43 step types, 172 Jinja filters, 1,664 playbooks",
            "FTS-indexed; one SQL query per agent question — no LLM where a lookup will do",
            "Probes refresh from the live appliance — compiler stays stable while the world changes",
            "Trust ladder (live_*, tested_pass → is_trusted=1) flags unverified data",
        ],
    },
    {
        "title": "Design 2 — Library-first compiler",
        "bullets": [
            "Pipeline: parser → resolver → validator → emitter",
            "Compiler is a Python library; CLI, MCP server, web app are thin wrappers",
            "Errors are data, not strings — every diagnostic has a code, line/col, suggested fix",
            "Round-trip (compile → decompile) is lossless modulo formatting",
        ],
    },
    {
        "title": "Design 3 — MCP as the agent contract",
        "bullets": [
            "20+ MCP tools: find_connector, get_op_schema, validate_yaml, compile_yaml,",
            "    render_jinja, run_op, get_run_env, resolve_picklist_value,",
            "    search_api_examples, synthesize_http_step, …",
            "The LLM never writes FSR JSON — it writes YAML and trusts the compiler",
            "LLM-agnostic by construction: structured I/O, externalized prompt, token discipline",
        ],
    },
    {
        "title": "The success ladder",
        "bullets": [
            "L1 Compile — structurally valid YAML",
            "L2 Static-resolve — connectors / ops / picklists / step-types / Jinja vars all exist",
            "L3 Dry-run — step args render against expected upstream context",
            "L4 Live single-step — step N actually executes on real FSR",
            "L5 Post-run assert — playbook produced the asked-for outcome",
            "Each rung is deterministic + an MCP tool; failure returns structured {ok, code, suggestions}",
        ],
        "diagram": "ladder",
    },
    {
        "title": "Recipes — from blank page to working playbook",
        "bullets": [
            "Recipe generators emit known-good playbooks for common archetypes",
            "Today: threat-feed ingestion, data ingestion (alerts/incidents)",
            "Roadmap: enrichment, triage, containment, approval-gated actions, orchestration",
            "Each recipe ships compile-clean, with TODO notes only where the user must intervene",
            "Layered ruleset validator enforces archetype-specific rules",
        ],
    },
    {
        "title": "HTTP virtual-connector — covering the long tail",
        "bullets": [
            "FortiSOAR ships a generic HTTP connector (10 ops, including http_paginate)",
            "We crawled 207,419 API examples across 6,927 third-party products",
            "search_api_examples + synthesize_http_step translate a catalog entry → http_request step",
            "Pre-fills method / path / auth / params from real examples — not LLM guesses",
            "Result: the agent reaches vendors we don't have a native connector for",
        ],
    },
    {
        "title": "Live verification loop",
        "bullets": [
            "validate_yaml catches shape errors before push",
            "render_jinja(template, from_pb_execution=<pk>) resolves Jinja against a real past run",
            "run_op executes one connector op live and caches observed output shape in SQLite",
            "dry_run_playbook (in build) imports → executes → cleans up on dev FSR",
            "Failures are evidence — if a demo breaks, that's a real bug to file",
        ],
    },
    {
        "title": "Demo 1 — vague ask to working playbook",
        "bullets": [
            "Prompt: \"Look up an IP on VirusTotal and update the alert\"",
            "Agent: find_connector → get_op_schema → emit YAML",
            "Then: validate_yaml → resolve_yaml (L2) → compile → push → run → assert",
            "Total time from prompt to passing run: minutes",
        ],
        "diagram": "flow",
    },
    {
        "title": "Demo 2 — the killer: triage and fix",
        "bullets": [
            "\"My playbook is broken. Figure out which one and fix it.\"",
            "Agent: list_recent_failed_runs → get_run_env(<pk>) → spots .records vs .data",
            "Pulls YAML, edits, validates, diffs, pushes, re-runs, watches it pass",
            "None of this loop is possible from the FSR Playbook Designer",
        ],
    },
    {
        "title": "What we've indexed — FortiSOAR knowledge",
        "subtitle_text": "Every number is a real row in a queryable index, refreshed from the live appliance.",
        "kind": "stats",
        "stats": [
            ("714", "Connectors"),
            ("6,773", "Operations"),
            ("26,093", "Op parameters"),
            ("43", "Step types"),
            ("172", "Jinja filters"),
            ("1,664", "Live playbooks"),
        ],
        "footer": "Trust ladder: 1,959 rows confirmed live + tested. The agent does SQL lookups before it touches an LLM.",
    },
    {
        "title": "What we've indexed — third-party APIs",
        "subtitle_text": "The HTTP virtual-connector corpus: thousands of vendors reachable without a native connector.",
        "kind": "stats",
        "stats": [
            ("6,927", "Products covered"),
            ("207,419", "API examples"),
            ("6,272", "Lifecycle records"),
            ("33,783", "Microsoft Graph alone"),
            ("13,818", "GitHub v3 REST"),
            ("FTS5", "Full-text indexed"),
        ],
        "footer": "search_api_examples + synthesize_http_step turn any of these into a working http_request step. Deterministic — no LLM in the translation.",
    },
    {
        "title": "What's already shipped",
        "bullets": [
            "714 connectors / 6,773 ops / 1,664 playbooks indexed in SQLite",
            "Compiler v1: parser → resolver → validator → emitter (round-trip lossless)",
            "20+ MCP tools, live-tested against a real FSR appliance",
            "E2E runner: 11/11 fixtures green; 4 LLM-driven storyboards in DEMO.md",
            "Recipe generators (feed + data ingest) + layered ruleset validator",
            "HTTP virtual-connector tools wired; inventory surface (fsrpb inventory)",
        ],
    },
    {
        "title": "What's next",
        "bullets": [
            "L2 prechecks (picklist resolvability, connector installation)",
            "resolve_yaml + variable-reachability ruleset — highest-ROI single check we don't have",
            "dry_run_playbook (L3) interactive stepper + assert_playbook_outcome (L5)",
            "LLM evaluation harness — proves \"any LLM works\" with measurements",
            "Recipe expansion: enrichment, triage, containment, HTTP virtual-connector recipes",
            "Inventory web dashboard + transcript capture for demo replay",
        ],
    },
    {
        "title": "Takeaways",
        "bullets": [
            "Authoring surface ≠ execution engine — keep them separate, win on both",
            "Agent + queryable spec + deterministic compiler + success ladder beats a visual designer",
            "YAML in git is the unfair advantage: review, audit, promote, roll back",
            "HTTP virtual-connector + 207K examples = thousands of vendors reachable today, no new code",
            "Q&A",
        ],
    },
]


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, prs):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_bullets(slide, bullets, prs):
    left = Inches(0.6)
    top = Inches(1.5)
    width = prs.slide_width - Inches(1.2)
    height = prs.slide_height - top - Inches(0.6)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + b
        run.font.name = "Helvetica"
        run.font.size = Pt(20)
        run.font.color.rgb = GREY


def _add_diagram_core(slide, prs):
    # three authoring surfaces -> IR -> compiler -> FSR JSON
    cx = prs.slide_width // 2
    top = Inches(4.6)
    box_w = Inches(2.0)
    box_h = Inches(0.7)
    gap = Inches(0.3)
    surfaces = ["Human (CLI)", "Agent (MCP)", "Visual editor"]
    total_w = box_w * 3 + gap * 2
    start = (prs.slide_width - total_w) // 2
    for i, label in enumerate(surfaces):
        x = start + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = NAVY
        shp.line.fill.background()
        tf = shp.text_frame
        tf.text = label
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.size = Pt(13); r.font.name = "Helvetica"
    # IR + compiler row
    row2_top = top + Inches(1.0)
    ir = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx - Inches(2.4), row2_top, Inches(2.0), box_h)
    ir.fill.solid(); ir.fill.fore_color.rgb = ACCENT; ir.line.fill.background()
    ir.text_frame.text = "YAML IR"
    for p in ir.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); r.font.size = Pt(13); r.font.bold = True
    cmp_ = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.4), row2_top, Inches(2.0), box_h)
    cmp_.fill.solid(); cmp_.fill.fore_color.rgb = NAVY; cmp_.line.fill.background()
    cmp_.text_frame.text = "Compiler → FSR JSON"
    for p in cmp_.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); r.font.size = Pt(13)


def _add_diagram_flow(slide, prs):
    steps = ["find_connector", "get_op_schema", "emit YAML", "validate", "compile", "push", "run"]
    top = Inches(4.8)
    box_w = Inches(1.25); box_h = Inches(0.6); gap = Inches(0.1)
    total = box_w * len(steps) + gap * (len(steps)-1)
    start = (prs.slide_width - total) // 2
    for i, label in enumerate(steps):
        x = start + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = NAVY if i % 2 == 0 else ACCENT
        shp.line.fill.background()
        shp.text_frame.text = label
        for p in shp.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs: r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); r.font.size = Pt(11); r.font.name = "Helvetica"


def _add_stats_grid(slide, prs, stats, subtitle_text=None, footer=None):
    if subtitle_text:
        _add_textbox(slide, Inches(0.6), Inches(1.25), prs.slide_width - Inches(1.2),
                     Inches(0.6), subtitle_text, size=14, color=GREY)
    cols = 3
    rows = (len(stats) + cols - 1) // cols
    cell_w = Inches(3.8); cell_h = Inches(1.4); gap = Inches(0.2)
    grid_w = cell_w * cols + gap * (cols - 1)
    start_x = (prs.slide_width - grid_w) // 2
    start_y = Inches(2.1)
    for i, (number, label) in enumerate(stats):
        r, c = divmod(i, cols)
        x = start_x + c * (cell_w + gap)
        y = start_y + r * (cell_h + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cell_w, cell_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = NAVY
        shp.line.fill.background()
        tf = shp.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.1)
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        run = p1.add_run(); run.text = number
        run.font.name = "Helvetica"; run.font.size = Pt(36); run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = label
        r2.font.name = "Helvetica"; r2.font.size = Pt(13)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xD4)
    if footer:
        _add_textbox(slide, Inches(0.6),
                     start_y + rows * (cell_h + gap) + Inches(0.1),
                     prs.slide_width - Inches(1.2), Inches(0.8),
                     footer, size=13, color=ACCENT, bold=True)


def _add_diagram_ladder(slide, prs):
    rungs = [
        ("L1", "Compile", "done"),
        ("L2", "Static-resolve", "partial"),
        ("L3", "Dry-run", "missing"),
        ("L4", "Live step", "half"),
        ("L5", "Post-run assert", "missing"),
    ]
    top = Inches(5.0)
    box_w = Inches(2.3); box_h = Inches(1.0); gap = Inches(0.15)
    total = box_w * len(rungs) + gap * (len(rungs) - 1)
    start = (prs.slide_width - total) // 2
    status_color = {
        "done": NAVY,
        "partial": RGBColor(0xC8, 0x86, 0x00),
        "half": RGBColor(0xC8, 0x86, 0x00),
        "missing": ACCENT,
    }
    for i, (rung, label, status) in enumerate(rungs):
        x = start + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = status_color[status]
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.text = f"{rung}\n{label}\n[{status}]"
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.size = Pt(12)
                r.font.name = "Helvetica"


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec in SLIDES:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, LIGHT)
        _add_title_bar(slide, prs)
        if spec.get("kind") == "title":
            _add_textbox(slide, Inches(0.6), Inches(2.6), prs.slide_width - Inches(1.2), Inches(1.4),
                         spec["title"], size=44, bold=True, color=NAVY)
            _add_textbox(slide, Inches(0.6), Inches(4.0), prs.slide_width - Inches(1.2), Inches(1.0),
                         spec["subtitle"], size=22, color=GREY)
            _add_textbox(slide, Inches(0.6), Inches(6.6), prs.slide_width - Inches(1.2), Inches(0.5),
                         "FSRPlaybookYaml — fsrpb", size=14, color=ACCENT, bold=True)
            continue
        _add_textbox(slide, Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(0.9),
                     spec["title"], size=30, bold=True, color=NAVY)
        if spec.get("kind") == "stats":
            _add_stats_grid(slide, prs, spec["stats"],
                            subtitle_text=spec.get("subtitle_text"),
                            footer=spec.get("footer"))
            continue
        if spec.get("bullets"):
            _add_bullets(slide, spec["bullets"], prs)
        if spec.get("diagram") == "core":
            _add_diagram_core(slide, prs)
        elif spec.get("diagram") == "flow":
            _add_diagram_flow(slide, prs)
        elif spec.get("diagram") == "ladder":
            _add_diagram_ladder(slide, prs)
    out = Path(__file__).resolve().parent.parent / "fsrpb_presentation.pptx"
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    build()
