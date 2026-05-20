"""Generate FSRPlaybookYaml presentation deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0F, 0x1E, 0x3D)
ACCENT = RGBColor(0xE8, 0x3E, 0x3E)
GREY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT = RGBColor(0xF4, 0xF4, 0xF6)

SLIDES = [
    {
        "title": "Agent-Authored Playbooks for FortiSOAR",
        "subtitle": "AI is the driver — but the deterministic cornerstones are why it works",
        "kind": "title",
    },
    {
        "title": "The pain we keep hitting",
        "bullets": [
            "Playbook Designer is click-heavy, opaque, and not diff-able",
            "Iteration loop is slow: edit → save → run → squint at run-log → repeat",
            "Field-shape bugs (.records vs .data, Jinja typos) only surface at runtime",
            "No version control, no PR review, no portability between environments",
            "Talented automation engineers spend their day clicking, not designing",
        ],
    },
    {
        "title": "The product principle",
        "bullets": [
            "\"Producing valid YAML is the start of the job, not the end.\"",
            "A playbook is useless if it doesn't run on a real FSR and produce the asked-for outcome",
            "So the system has to prove playbooks work, not just compile",
            "Everything that follows serves that one principle",
        ],
    },
    {
        "title": "Why AI is the right driver here",
        "bullets": [
            "Playbook authoring is structured + bounded — exactly where modern LLMs shine",
            "The hard parts are knowing which connector / op / picklist exists, not invention",
            "Natural language → YAML is the right abstraction step; LLMs skip the click-through",
            "Agentic tool use means the model asks the system instead of guessing",
            "Bottleneck moved: not \"can the model write YAML\" but \"can it verify what it wrote\"",
        ],
    },
    {
        "title": "But this is more than \"just AI\"",
        "bullets": [
            "An LLM alone hallucinates connectors, invents picklist values, mistypes vars.steps paths",
            "Without verification it's a faster way to ship broken playbooks",
            "The cornerstones below are deterministic, queryable, and ground-truth — they make AI safe",
            "Each one stands alone as a tool a human can run; together they close the loop on AI",
            "If you remove the AI tomorrow, every cornerstone is still a useful authoring + audit tool",
        ],
    },
    {
        "title": "The cornerstones (six pillars under the AI layer)",
        "bullets": [
            "1.  SQLite reference store — connectors, ops, picklists, step types, Jinja, recipes",
            "2.  Library-first compiler — parser → resolver → validator → emitter, errors as data",
            "3.  Render-path validator — simulate playbooks offline, surface data-access bugs",
            "4.  Live-FSR probes — capture real runtime semantics into reusable fixtures",
            "5.  Step-param audit — corpus-driven param truth, surfaces resolver gaps",
            "6.  MCP toolbelt — every cornerstone is callable by the agent or a human",
        ],
    },
    {
        "title": "The core idea",
        "bullets": [
            "One simplified YAML IR. Three authoring surfaces (CLI, agent-via-MCP, future visual editor)",
            "All three compile through the same pipeline to native FSR WorkflowCollection JSON",
            "FortiSOAR stays the execution engine — we never replace it, we author into it",
            "Same IR across surfaces → output is interoperable",
        ],
        "diagram": "core",
    },
    {
        "title": "Why YAML is the right IR",
        "bullets": [
            "Human-readable, diff-able, review-able in a PR",
            "Trivially round-trippable: pull from FSR → YAML → edit → push",
            "Easy for an LLM to emit and reason about",
            "No vendor JSON ceremony (UUIDs, port coords, layout metadata) leaks into authoring",
        ],
    },
    {
        "title": "Design 1 — SQLite-first reference store",
        "bullets": [
            "714 connectors, 6,773 ops, 26K params, 43 step types, 172 Jinja filters, 1,664 playbooks",
            "FTS-indexed; one SQL query per agent question — no LLM where a lookup will do",
            "Probes refresh from the live appliance — compiler stays stable while the world changes",
            "Trust ladder (live_*, tested_pass → is_trusted=1) flags unverified data",
        ],
    },
    {
        "title": "Design 2 — Library-first compiler",
        "bullets": [
            "Pipeline: parser → resolver → validator → emitter",
            "Compiler is a Python library; CLI, MCP server, web app are thin wrappers",
            "Errors are data, not strings — every diagnostic has a code, line/col, suggested fix",
            "Round-trip (compile → decompile) is lossless modulo formatting",
        ],
    },
    {
        "title": "Cornerstone 3 — Render-path validator",
        "bullets": [
            "Simulates the playbook offline before any push: every step's rendered args + output shape",
            "Per-step provenance: mock_result · computed · live_run · default_empty",
            "Catches what validate_yaml can't: vars.steps.X.Y typos, missing keys, type drift",
            "Static extractor finds every vars.…, picklist(), Jinja reference per step",
            "Heuristic checks (C1-C4): unreachable refs, missing keys, required-empty, picklist drift",
            "Surfaces close-key suggestions deterministically — no LLM in the diagnostic itself",
        ],
    },
    {
        "title": "Cornerstone 4 — Live-FSR probes ground simulator truth",
        "bullets": [
            "Probes synthesize tiny playbooks, push, fire, capture full run env into JSON fixtures",
            "9 scenarios so far: for_each (sequential / parallel / break / empty / condition), workflow_reference",
            "Catches non-obvious truths: vars.steps.<for_each> is list-of-dicts, not last-value",
            "Pinned in agent prompt as facts the LLM must not invent",
            "Same fixtures back parametrized parity tests — drift on either side surfaces immediately",
            "Re-runnable per scenario; new constructs become probes, not guesses",
        ],
    },
    {
        "title": "Cornerstone 5 — Step-param audit (corpus-driven truth)",
        "bullets": [
            "fsrpb dump-step-params writes one Markdown report per step type",
            "Combines: resolver allowlist + 7,442 corpus rows of real arg keys + frequency",
            "Highlights gaps: resolver doesn't accept a key the corpus uses 305× (e.g. FindRecords.query)",
            "Already surfaced 5+ resolver allowlist holes; each gap is a one-line fix",
            "Lives in docs/step_params/ — agent greps these instead of trusting training intuition",
            "Re-generated on demand; stays aligned as FSR evolves",
        ],
    },
    {
        "title": "Cornerstone 6 — MCP as the agent contract",
        "bullets": [
            "30+ MCP tools — every cornerstone is callable by name with structured I/O",
            "find_connector · get_op_schema · validate_yaml · compile_yaml · render_jinja",
            "analyze_playbook · suggest_fix_for_diagnostic · step_test · step_through_playbook",
            "run_op · get_run_env · resolve_picklist_value · find_step_examples",
            "search_api_examples · synthesize_http_step · find_recipe",
            "LLM-agnostic by construction: any model that does tool use can drive this",
        ],
    },
    {
        "title": "The success ladder",
        "bullets": [
            "Compiles — structurally valid YAML + every Jinja var reachable from a prior step",
            "Runs — executes against real FSR (static-resolve + dry-run roll up here)",
            "Works — post-run assertions hold (record exists, field equals, count > N)",
            "Each rung is deterministic + an MCP tool; failure returns structured {ok, code, suggestions}",
        ],
        "diagram": "ladder",
    },
    {
        "title": "The agent prompt — workflow gates that close the loop",
        "bullets": [
            "Hard rules pin the YAML shape — no negotiation on canonical form",
            "Required workflow: recipe lookup → schema lookup → draft → validate → analyze → done",
            "Mandatory analyze gate: do NOT declare done until error_count == 0",
            "Per-diagnostic-kind fix strategies baked into the prompt — no prose freelancing",
            "FSR runtime semantics block: live-verified shapes the LLM must not invent",
            "Tool error contract: failures return {ok:false, code, suggestions} — agent retries structurally",
        ],
    },
    {
        "title": "Two distinct agent flows",
        "bullets": [
            "Authoring (new playbook): intent → recipes → schemas → YAML → validate → analyze → push",
            "    \"Build me a playbook to enrich every alert with VirusTotal\"",
            "Troubleshooting (existing): pull → analyze → suggest_fix → patch → re-analyze → diff → push",
            "    \"My playbook is broken — figure it out\"",
            "Same MCP tools, same cornerstones — what changes is the entry point + the gate order",
            "If the agent skips the analyze gate, the runtime catches it. The system is the safety net.",
        ],
        "diagram": "agent_flow",
    },
    {
        "title": "The visual editor closes the loop for humans",
        "bullets": [
            "Same cornerstones, rendered as canvas chrome:",
            "    • Per-node red/amber badge keyed off worst diagnostic for that step",
            "    • Diagnostics drawer lists every render-path issue with kind + location + suggestion",
            "    • Suggest-fix button per row → before/after diff → Apply → auto re-analyze",
            "    • Click step_id → canvas jumps to the offending node, inspector opens",
            "    • Node colors mirror the simulator's family tagging — no hidden context",
            "Every action goes through the same MCP tools the agent uses. One source of truth.",
        ],
    },
    {
        "title": "The full MCP toolbelt — 49 tools the agent can call",
        "kind": "tool_catalog",
        "subtitle_text": "Same tools serve chat, CLI, and visual editor. Every box below is one MCP call away.",
        "groups": [
            ("Discover", [
                "find_connector", "find_operation", "find_step_type",
                "find_step_examples", "find_operation_example",
                "find_recipe", "find_step_recipe",
                "find_jinja_filter", "find_jinja_pattern",
                "find_jinja_example", "get_filter_examples",
                "get_op_schema", "get_step_type", "get_picklist",
                "list_picklists", "picklist_for_field",
                "search_playbooks", "search_api_examples",
                "list_tags", "list_configured_connectors",
                "list_connector_configurations",
                "get_connector_icon", "get_connector_source",
            ]),
            ("Build", [
                "generate_recipe", "synthesize_http_step",
                "resolve_picklist_value",
            ]),
            ("Validate (offline)", [
                "validate_yaml", "compile_yaml", "resolve_yaml",
                "render_jinja", "analyze_playbook",
                "suggest_fix_for_diagnostic",
                "step_through_playbook", "step_test",
                "precheck_picklist_value",
                "precheck_connector_installed",
                "healthcheck_connector",
            ]),
            ("Live", [
                "push_playbook", "run_op", "run_playbook",
                "dry_run_playbook", "search_module_records",
                "test_find_record",
            ]),
            ("Diagnose", [
                "list_recent_failed_runs", "list_playbook_runs",
                "get_run_env",
                "diagnose_yaml_against_pb_execution",
                "assert_playbook_outcome", "verification_status",
                "review_chat_session", "review_recent_thumbs_down",
            ]),
        ],
    },
    {
        "title": "Recipes — from blank page to working playbook",
        "bullets": [
            "Recipe generators emit known-good playbooks for common archetypes",
            "Today: threat-feed ingestion, data ingestion (alerts/incidents)",
            "Roadmap: enrichment, triage, containment, approval-gated actions, orchestration",
            "Each recipe ships compile-clean, with TODO notes only where the user must intervene",
            "Layered ruleset validator enforces archetype-specific rules",
        ],
    },
    {
        "title": "HTTP virtual-connector — covering the long tail",
        "bullets": [
            "FortiSOAR ships a generic HTTP connector (10 ops, including http_paginate)",
            "We crawled 207,419 API examples across 6,927 third-party products",
            "search_api_examples + synthesize_http_step translate a catalog entry → http_request step",
            "Pre-fills method / path / auth / params from real examples — not LLM guesses",
            "Result: the agent reaches vendors we don't have a native connector for",
        ],
    },
    {
        "title": "Roadmap — vendor doc ingest (zero-day SaaS support)",
        "bullets": [
            "Today: 207K examples crawled from API marketplaces — already covers most major vendors",
            "Next: ingest vendors directly from their published docs the day they ship them",
            "    • OpenAPI/Swagger specs — most modern SaaS vendors publish one; full op + schema graph",
            "    • Postman collections — covers vendors that ship Postman instead of OpenAPI",
            "    • AsyncAPI / GraphQL schemas — adds streaming + graph endpoints to the catalog",
            "Auth pattern library: cluster recurring auth flows (Bearer / OAuth2 / HMAC / JWT)",
            "    so the agent picks a named pattern instead of re-deriving from one example",
            "Per-vendor recipe scaffolds: spec metadata → ready-made enrich / contain / triage starters",
            "Outcome: a new vendor goes live in fsrpb the same day they publish their API docs",
        ],
    },
    {
        "title": "Live verification loop",
        "bullets": [
            "validate_yaml catches shape errors before push",
            "render_jinja(template, from_pb_execution=<pk>) resolves Jinja against a real past run",
            "run_op executes one connector op live and caches observed output shape in SQLite",
            "dry_run_playbook (in build) imports → executes → cleans up on dev FSR",
            "Failures are evidence — if a demo breaks, that's a real bug to file",
        ],
    },
    {
        "title": "Demo 1 — vague ask to working playbook",
        "bullets": [
            "Prompt: \"Look up an IP on VirusTotal and update the alert\"",
            "Agent: find_recipe → find_connector → get_op_schema → draft YAML",
            "Gates: validate_yaml → analyze_playbook → (only if 0 errors) compile → push",
            "Render-path validator catches a vars.steps typo in seconds, before any push",
            "Total time prompt-to-passing-run: minutes — and every step is auditable in chat",
        ],
        "diagram": "flow",
    },
    {
        "title": "Demo 2 — the killer: triage and fix",
        "bullets": [
            "\"My playbook is broken. Figure out which one and fix it.\"",
            "Agent: pull → analyze_playbook → diagnostics grouped by step + severity",
            "C2 missing_key fires: \"data.statuss\" — close-key suggestion: \"status\"",
            "suggest_fix_for_diagnostic returns a structured patch — agent applies, re-analyzes",
            "Diff shown to user, push on confirm — none of this is possible in FSR Designer",
        ],
    },
    {
        "title": "What we've indexed — FortiSOAR knowledge",
        "subtitle_text": "Every number is a real row in a queryable index, refreshed from the live appliance.",
        "kind": "stats",
        "stats": [
            ("714", "Connectors"),
            ("6,773", "Operations"),
            ("26,093", "Op parameters"),
            ("43", "Step types"),
            ("172", "Jinja filters"),
            ("1,664", "Live playbooks"),
        ],
        "footer": "Trust ladder: 1,959 rows confirmed live + tested. The agent does SQL lookups before it touches an LLM.",
    },
    {
        "title": "What we've indexed — third-party APIs",
        "subtitle_text": "The HTTP virtual-connector corpus: thousands of vendors reachable without a native connector.",
        "kind": "stats",
        "stats": [
            ("6,927", "Products covered"),
            ("207,419", "API examples"),
            ("6,272", "Lifecycle records"),
            ("33,783", "Microsoft Graph alone"),
            ("13,818", "GitHub v3 REST"),
            ("FTS5", "Full-text indexed"),
        ],
        "footer": "search_api_examples + synthesize_http_step turn any of these into a working http_request step. Deterministic — no LLM in the translation.",
    },
    {
        "title": "What's already shipped",
        "bullets": [
            "714 connectors / 6,773 ops / 7,442 corpus playbook-steps indexed in SQLite",
            "Compiler v1: parser → resolver → validator → emitter (round-trip lossless)",
            "Render-path validator: simulator + 4 heuristic checks (C1-C4) + suggest_fix_for_diagnostic",
            "9 live-FSR probe scenarios + 22 fixture-backed parity tests pinning runtime semantics",
            "Step-param audit (docs/step_params/) — 23 reports, surfacing real resolver gaps",
            "Visual editor diagnostics: per-node badges + drawer + suggest-fix Apply loop",
            "30+ MCP tools; agent prompt with mandatory analyze gate; 407 Python + 214 JS tests green",
            "fsrpb CLI: validate · analyze · compile · push · diff · dump-step-params",
        ],
    },
    {
        "title": "What's next",
        "bullets": [
            "Phase 5 analyzer checks: type mismatch, index-into-non-list, dead-step detection",
            "Recursive output_shape — catch typos in nested paths (vars.steps.X.data.foo)",
            "Resolver gap close-out — the dump-step-params reports already list each by row count",
            "Inspector \"Render path\" tab: rendered args + simulated output + consumed paths inline",
            "Bulk fix: same typo across N steps → one click patches all",
            "LLM evaluation harness — measure \"any model works\" with repeatable scenarios",
            "Vendor doc ingest (see preceding slide) — zero-day SaaS support",
            "Recipe expansion: enrichment, triage, containment, approval-gated workflows",
        ],
    },
    {
        "title": "Takeaways",
        "bullets": [
            "AI is the right driver, but only because the cornerstones underneath are deterministic",
            "Every cornerstone (compiler, validator, probes, audit) stands alone as a useful tool",
            "The agent prompt's analyze gate is the system's seatbelt — verifiable, not aspirational",
            "Same MCP tools serve chat, CLI, and visual editor — one source of truth",
            "YAML in git + analyze-before-push = playbooks that can be reviewed, promoted, rolled back",
            "If you swap the LLM tomorrow, everything still works. That's the value of the chassis.",
        ],
    },
]


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, prs):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_bullets(slide, bullets, prs):
    left = Inches(0.6)
    top = Inches(1.5)
    width = prs.slide_width - Inches(1.2)
    height = prs.slide_height - top - Inches(0.6)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + b
        run.font.name = "Helvetica"
        run.font.size = Pt(20)
        run.font.color.rgb = GREY


def _add_diagram_core(slide, prs):
    # three authoring surfaces -> IR -> compiler -> FSR JSON
    cx = prs.slide_width // 2
    top = Inches(4.6)
    box_w = Inches(2.0)
    box_h = Inches(0.7)
    gap = Inches(0.3)
    surfaces = ["Human (CLI)", "Agent (MCP)", "Visual editor"]
    total_w = box_w * 3 + gap * 2
    start = (prs.slide_width - total_w) // 2
    for i, label in enumerate(surfaces):
        x = start + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = NAVY
        shp.line.fill.background()
        tf = shp.text_frame
        tf.text = label
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.size = Pt(13); r.font.name = "Helvetica"
    # IR + compiler row
    row2_top = top + Inches(1.0)
    ir = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx - Inches(2.4), row2_top, Inches(2.0), box_h)
    ir.fill.solid(); ir.fill.fore_color.rgb = ACCENT; ir.line.fill.background()
    ir.text_frame.text = "YAML IR"
    for p in ir.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); r.font.size = Pt(13); r.font.bold = True
    cmp_ = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.4), row2_top, Inches(2.0), box_h)
    cmp_.fill.solid(); cmp_.fill.fore_color.rgb = NAVY; cmp_.line.fill.background()
    cmp_.text_frame.text = "Compiler → FSR JSON"
    for p in cmp_.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); r.font.size = Pt(13)


def _add_diagram_flow(slide, prs):
    steps = ["find_connector", "get_op_schema", "emit YAML", "validate", "compile", "push", "run"]
    top = Inches(4.8)
    box_w = Inches(1.25); box_h = Inches(0.6); gap = Inches(0.1)
    total = box_w * len(steps) + gap * (len(steps)-1)
    start = (prs.slide_width - total) // 2
    for i, label in enumerate(steps):
        x = start + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = NAVY if i % 2 == 0 else ACCENT
        shp.line.fill.background()
        shp.text_frame.text = label
        for p in shp.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs: r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); r.font.size = Pt(11); r.font.name = "Helvetica"


def _add_stats_grid(slide, prs, stats, subtitle_text=None, footer=None):
    if subtitle_text:
        _add_textbox(slide, Inches(0.6), Inches(1.25), prs.slide_width - Inches(1.2),
                     Inches(0.6), subtitle_text, size=14, color=GREY)
    cols = 3
    rows = (len(stats) + cols - 1) // cols
    cell_w = Inches(3.8); cell_h = Inches(1.4); gap = Inches(0.2)
    grid_w = cell_w * cols + gap * (cols - 1)
    start_x = (prs.slide_width - grid_w) // 2
    start_y = Inches(2.1)
    for i, (number, label) in enumerate(stats):
        r, c = divmod(i, cols)
        x = start_x + c * (cell_w + gap)
        y = start_y + r * (cell_h + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cell_w, cell_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = NAVY
        shp.line.fill.background()
        tf = shp.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.1)
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        run = p1.add_run(); run.text = number
        run.font.name = "Helvetica"; run.font.size = Pt(36); run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = label
        r2.font.name = "Helvetica"; r2.font.size = Pt(13)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xD4)
    if footer:
        _add_textbox(slide, Inches(0.6),
                     start_y + rows * (cell_h + gap) + Inches(0.1),
                     prs.slide_width - Inches(1.2), Inches(0.8),
                     footer, size=13, color=ACCENT, bold=True)


def _add_diagram_ladder(slide, prs):
    rungs = [
        ("", "Compiles", "done"),
        ("", "Runs", "partial"),
        ("", "Works", "missing"),
    ]
    top = Inches(5.0)
    box_w = Inches(2.3); box_h = Inches(1.0); gap = Inches(0.15)
    total = box_w * len(rungs) + gap * (len(rungs) - 1)
    start = (prs.slide_width - total) // 2
    status_color = {
        "done": NAVY,
        "partial": RGBColor(0xC8, 0x86, 0x00),
        "half": RGBColor(0xC8, 0x86, 0x00),
        "missing": ACCENT,
    }
    for i, (rung, label, status) in enumerate(rungs):
        x = start + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = status_color[status]
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.text = f"{rung}\n{label}\n[{status}]"
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.size = Pt(12)
                r.font.name = "Helvetica"


def _add_tool_catalog(slide, prs, groups, subtitle_text=None):
    """5-column layout, one column per agent-flow phase. Each column
    renders the phase label as a colored header and the tools below
    as a tight monospace list. Sizes auto-scale to fit the longest
    column without overflow."""
    if subtitle_text:
        _add_textbox(slide, Inches(0.6), Inches(1.15),
                     prs.slide_width - Inches(1.2), Inches(0.5),
                     subtitle_text, size=13, color=GREY)

    n_cols = len(groups)
    avail_w = prs.slide_width - Inches(1.2)
    gap = Inches(0.15)
    col_w = (avail_w - gap * (n_cols - 1)) / n_cols
    top = Inches(1.8)
    col_h = prs.slide_height - top - Inches(0.5)
    header_h = Inches(0.45)

    palette = [NAVY, ACCENT,
               RGBColor(0x0E, 0x6B, 0x4F),
               RGBColor(0xC8, 0x86, 0x00),
               RGBColor(0x5A, 0x2D, 0x82)]

    for i, (label, tools) in enumerate(groups):
        x = Inches(0.6) + i * (col_w + gap)
        # Header bar
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, col_w, header_h)
        hdr.fill.solid(); hdr.fill.fore_color.rgb = palette[i % len(palette)]
        hdr.line.fill.background()
        tf = hdr.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = f"{label}  ({len(tools)})"
        run.font.name = "Helvetica"; run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Tool list body — mono font, tight spacing.
        body_top = top + header_h + Inches(0.05)
        body = slide.shapes.add_textbox(x, body_top, col_w,
                                        col_h - header_h - Inches(0.05))
        btf = body.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.1); btf.margin_right = Inches(0.05)
        btf.margin_top = Inches(0.05); btf.margin_bottom = Inches(0.05)
        # Auto-size: scale font down for dense columns so they fit
        # the column height without spill.
        n = len(tools)
        size = 10 if n <= 12 else 9 if n <= 18 else 8
        for j, name in enumerate(tools):
            para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            r = para.add_run()
            r.text = name
            r.font.name = "Menlo"
            r.font.size = Pt(size)
            r.font.color.rgb = NAVY


def _add_diagram_agent_flow(slide, prs):
    """Two side-by-side agent loops: authoring (left) and
    troubleshooting (right). Each is a horizontal flow of stations."""
    top = Inches(4.4)
    box_h = Inches(0.55)
    half_w = (prs.slide_width - Inches(1.2)) // 2 - Inches(0.15)
    # Authoring lane
    auth_x = Inches(0.6)
    auth_label = slide.shapes.add_textbox(auth_x, top - Inches(0.45),
                                          half_w, Inches(0.35))
    auth_label.text_frame.text = "Authoring (new)"
    for p in auth_label.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = NAVY
            r.font.name = "Helvetica"

    auth_steps = ["intent", "find_recipe", "find_operation",
                  "draft YAML", "validate", "analyze", "push"]
    sw = (half_w - Inches(0.05) * (len(auth_steps) - 1)) / len(auth_steps)
    for i, label in enumerate(auth_steps):
        x = auth_x + i * (sw + Inches(0.05))
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, top, sw, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = NAVY
        shp.line.fill.background()
        shp.text_frame.text = label
        for p in shp.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.size = Pt(9); r.font.name = "Helvetica"

    # Troubleshooting lane
    tr_x = auth_x + half_w + Inches(0.3)
    tr_label = slide.shapes.add_textbox(tr_x, top - Inches(0.45),
                                        half_w, Inches(0.35))
    tr_label.text_frame.text = "Troubleshooting (existing)"
    for p in tr_label.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = ACCENT
            r.font.name = "Helvetica"

    tr_steps = ["pull", "analyze", "diagnostics",
                "suggest_fix", "patch", "re-analyze", "diff/push"]
    sw2 = (half_w - Inches(0.05) * (len(tr_steps) - 1)) / len(tr_steps)
    for i, label in enumerate(tr_steps):
        x = tr_x + i * (sw2 + Inches(0.05))
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, top, sw2, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = ACCENT
        shp.line.fill.background()
        shp.text_frame.text = label
        for p in shp.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.size = Pt(9); r.font.name = "Helvetica"

    # Footer note tying both lanes back to the cornerstones.
    _add_textbox(slide, Inches(0.6), top + Inches(1.3),
                 prs.slide_width - Inches(1.2), Inches(0.4),
                 "Both flows share the same MCP tools. The analyze gate "
                 "is the seatbelt — it can't be bypassed by either flow.",
                 size=12, color=GREY, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec in SLIDES:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, LIGHT)
        _add_title_bar(slide, prs)
        if spec.get("kind") == "title":
            _add_textbox(slide, Inches(0.6), Inches(2.6), prs.slide_width - Inches(1.2), Inches(1.4),
                         spec["title"], size=44, bold=True, color=NAVY)
            _add_textbox(slide, Inches(0.6), Inches(4.0), prs.slide_width - Inches(1.2), Inches(1.0),
                         spec["subtitle"], size=22, color=GREY)
            _add_textbox(slide, Inches(0.6), Inches(6.6), prs.slide_width - Inches(1.2), Inches(0.5),
                         "FSRPlaybookYaml — fsrpb", size=14, color=ACCENT, bold=True)
            continue
        _add_textbox(slide, Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(0.9),
                     spec["title"], size=30, bold=True, color=NAVY)
        if spec.get("kind") == "stats":
            _add_stats_grid(slide, prs, spec["stats"],
                            subtitle_text=spec.get("subtitle_text"),
                            footer=spec.get("footer"))
            continue
        if spec.get("kind") == "tool_catalog":
            _add_tool_catalog(slide, prs, spec["groups"],
                              subtitle_text=spec.get("subtitle_text"))
            continue
        if spec.get("bullets"):
            _add_bullets(slide, spec["bullets"], prs)
        if spec.get("diagram") == "core":
            _add_diagram_core(slide, prs)
        elif spec.get("diagram") == "flow":
            _add_diagram_flow(slide, prs)
        elif spec.get("diagram") == "ladder":
            _add_diagram_ladder(slide, prs)
        elif spec.get("diagram") == "agent_flow":
            _add_diagram_agent_flow(slide, prs)
    out = Path(__file__).resolve().parent.parent / "fsrpb_presentation.pptx"
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    build()
